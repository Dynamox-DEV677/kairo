"""
Kairo · Pitch Deck v2 — 12-slide cover-to-CTA presentation.

Strict monochrome (black · deep purple · white). Editorial / brutalist
language matching the landing page and recruitment poster.

Run:
    pip install python-pptx
    python tools/make_kairo_pitch_deck.py

Output:  Kairo_Pitch_Deck_v2.pptx (repo root)
"""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ─── palette ───────────────────────────────────────────────────────────────
BG          = RGBColor(0x06, 0x06, 0x0A)
PANEL       = RGBColor(0x0E, 0x0E, 0x16)
PANEL2      = RGBColor(0x14, 0x14, 0x1F)
BORDER      = RGBColor(0x22, 0x22, 0x2E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DIM    = RGBColor(0xC1, 0xC1, 0xC8)
TEXT_FAINT  = RGBColor(0x8A, 0x8A, 0x96)
TEXT_VERY   = RGBColor(0x5A, 0x5A, 0x66)
PURPLE_LITE = RGBColor(0xE9, 0xD5, 0xFF)
PURPLE_SOFT = RGBColor(0xC4, 0xB5, 0xFD)
PURPLE      = RGBColor(0xA7, 0x8B, 0xFA)
PURPLE_HI   = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_DEEP = RGBColor(0x5B, 0x21, 0xB6)
PURPLE_INK  = RGBColor(0x3B, 0x07, 0x64)

# 16:9 slide dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ─── helpers ───────────────────────────────────────────────────────────────
def set_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()

def add_rect(slide, x, y, w, h, rgb=PANEL, line_rgb=None, corner=0.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if corner:
        shape.adjustments[0] = corner
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape

def add_text(slide, x, y, w, h, text, *, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Inter',
             italic=False, letter_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top  = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *, size=14, bullet_color=PURPLE_SOFT,
                text_color=TEXT_DIM, bold_label=True, line_gap=8, font='Inter'):
    """items: list of (label, body)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top  = tf.margin_bottom = Pt(0)
    for i, (label, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(line_gap)
        r0 = p.add_run()
        r0.text = '•  '
        r0.font.name = font
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.color.rgb = bullet_color
        if label:
            r1 = p.add_run()
            r1.text = label
            r1.font.name = font
            r1.font.size = Pt(size)
            r1.font.bold = bold_label
            r1.font.color.rgb = WHITE
            r2 = p.add_run()
            r2.text = (' — ' + body) if body else ''
            r2.font.name = font
            r2.font.size = Pt(size)
            r2.font.color.rgb = text_color
        else:
            r2 = p.add_run()
            r2.text = body
            r2.font.name = font
            r2.font.size = Pt(size)
            r2.font.color.rgb = text_color
    return tb

def add_background(slide, rgb=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb
    bg.line.fill.background()

def add_corner_glow(slide, alpha_top=0x18, alpha_bot=0x12):
    g1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
    g1.fill.solid()
    g1.fill.fore_color.rgb = PURPLE_DEEP
    g1.line.fill.background()
    _set_alpha(g1, alpha_top)
    g2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(5), Inches(6), Inches(6))
    g2.fill.solid()
    g2.fill.fore_color.rgb = PURPLE_HI
    g2.line.fill.background()
    _set_alpha(g2, alpha_bot)

def _set_alpha(shape, alpha_byte):
    sp = shape.fill._xPr.find(qn('a:solidFill'))
    if sp is None: return
    clr = sp.find(qn('a:srgbClr'))
    if clr is None: return
    for a in clr.findall(qn('a:alpha')):
        clr.remove(a)
    a = etree.SubElement(clr, qn('a:alpha'))
    val = int(alpha_byte / 255 * 100000)
    a.set('val', str(val))

def add_chip(slide, x, y, w, h, text, rgb_bg=PANEL2, rgb_border=BORDER,
             rgb_text=PURPLE_SOFT, size=10):
    shape = add_rect(slide, x, y, w, h, rgb=rgb_bg, line_rgb=rgb_border, corner=0.5)
    tf = shape.text_frame
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top  = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper()
    r.font.name = 'Inter'
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = rgb_text
    return shape

def add_eyebrow(slide, x, y, num, label):
    add_chip(slide, x, y, Inches(2.4), Inches(0.36), label, rgb_text=PURPLE_SOFT, size=11)
    add_text(slide, x, y + Inches(0.55), Inches(2), Inches(0.6),
             num + '.', size=44, bold=True, color=PURPLE, letter_spacing=-2)

def add_purple_rule(slide, x, y, w, thick=Pt(3)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, thick)
    line.fill.solid(); line.fill.fore_color.rgb = PURPLE
    line.line.fill.background()

def add_footer(slide, page, total):
    add_text(slide, Inches(0.55), Inches(7.08), Inches(6), Inches(0.3),
             'KAIRO · ACCELERATE YOUR ACADEMICS',
             size=9, color=TEXT_VERY)
    add_text(slide, Inches(7.0), Inches(7.08), Inches(5.8), Inches(0.3),
             f'PITCH №01  ·  {page} / {total}',
             size=9, color=TEXT_VERY, align=PP_ALIGN.RIGHT)


# ─── slide builders ────────────────────────────────────────────────────────
def slide_blank(prs):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    add_corner_glow(slide)
    return slide

def slide_cover(prs, total):
    slide = slide_blank(prs)
    # Top issue strip
    add_text(slide, Inches(0.6), Inches(0.5), Inches(6), Inches(0.4),
             'KAIRO · EDU-OS', size=12, bold=True, color=WHITE)
    add_text(slide, Inches(7.3), Inches(0.5), Inches(5.5), Inches(0.4),
             'ISSUE №01  ·  MAY 2026  ·  CHENNAI', size=12, color=TEXT_FAINT,
             align=PP_ALIGN.RIGHT)
    # Centre divider rule
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), Inches(1.0), Inches(12.13), Pt(1.2))
    line.fill.solid(); line.fill.fore_color.rgb = BORDER
    line.line.fill.background()
    # Massive type
    add_text(slide, Inches(0.6), Inches(2.0), Inches(12), Inches(1.5),
             'KAIRO.', size=140, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, Inches(0.6), Inches(3.9), Inches(12), Inches(1.0),
             'An AI that remembers your mind.',
             size=34, color=PURPLE_SOFT, italic=True)
    # Purple bar
    add_purple_rule(slide, Inches(0.6), Inches(5.0), Inches(2.2), thick=Pt(3))
    # Tag block
    add_text(slide, Inches(0.6), Inches(5.3), Inches(12), Inches(0.55),
             'A learning operating system for Class 9–12 India.',
             size=18, color=TEXT_DIM)
    add_text(slide, Inches(0.6), Inches(5.85), Inches(12), Inches(0.55),
             'CBSE · ICSE · State boards.',
             size=13, color=TEXT_FAINT)
    # Bottom credit
    add_text(slide, Inches(0.6), Inches(6.7), Inches(12), Inches(0.4),
             'Built by Darshan — 13, Chennai · 2026',
             size=11, color=TEXT_VERY)
    add_footer(slide, 1, total)
    return slide

def slide_section(prs, num, eyebrow, headline, lede, page, total):
    slide = slide_blank(prs)
    add_eyebrow(slide, Inches(0.55), Inches(0.5), num, eyebrow)
    add_text(slide, Inches(0.55), Inches(2.4), Inches(12), Inches(1.7),
             headline, size=52, bold=True, color=WHITE)
    add_purple_rule(slide, Inches(0.6), Inches(4.3), Inches(2.2), thick=Pt(3))
    add_text(slide, Inches(0.55), Inches(4.7), Inches(11.5), Inches(2.0),
             lede, size=17, color=TEXT_DIM)
    add_footer(slide, page, total)
    return slide

def slide_bullets(prs, num, eyebrow, headline, items, page, total, foot_note=None):
    slide = slide_blank(prs)
    add_eyebrow(slide, Inches(0.55), Inches(0.5), num, eyebrow)
    add_text(slide, Inches(0.55), Inches(2.2), Inches(12), Inches(1.2),
             headline, size=40, bold=True, color=WHITE)
    add_purple_rule(slide, Inches(0.6), Inches(3.45), Inches(1.8), thick=Pt(3))
    add_bullets(slide, Inches(0.6), Inches(3.85), Inches(12), Inches(3.2),
                items, size=14.5, line_gap=12)
    if foot_note:
        add_text(slide, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
                 foot_note, size=10, color=TEXT_VERY)
    add_footer(slide, page, total)
    return slide

def slide_stat_grid(prs, num, eyebrow, headline, stats, page, total, note=None):
    """stats: list of (big_number, label) — exactly 4."""
    slide = slide_blank(prs)
    add_eyebrow(slide, Inches(0.55), Inches(0.5), num, eyebrow)
    add_text(slide, Inches(0.55), Inches(2.2), Inches(12), Inches(1.2),
             headline, size=40, bold=True, color=WHITE)
    add_purple_rule(slide, Inches(0.6), Inches(3.45), Inches(1.8), thick=Pt(3))

    grid_y = Inches(3.95)
    col_w  = Inches(2.95)
    gap    = Inches(0.1)
    for i, (num_str, label) in enumerate(stats):
        x = Inches(0.55) + i * (col_w + gap)
        # Card
        add_rect(slide, x, grid_y, col_w, Inches(2.45), rgb=PANEL, line_rgb=BORDER, corner=0.05)
        add_text(slide, x + Inches(0.25), grid_y + Inches(0.4), col_w - Inches(0.5), Inches(1.2),
                 num_str, size=48, bold=True, color=WHITE, letter_spacing=-1)
        # purple chip below big number
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x + Inches(0.25), grid_y + Inches(1.6),
                                      Inches(0.6), Pt(2))
        line.fill.solid(); line.fill.fore_color.rgb = PURPLE
        line.line.fill.background()
        add_text(slide, x + Inches(0.25), grid_y + Inches(1.75), col_w - Inches(0.5), Inches(0.6),
                 label, size=11, color=PURPLE_SOFT, font='Inter')
    if note:
        add_text(slide, Inches(0.55), Inches(6.75), Inches(12), Inches(0.4),
                 note, size=11, color=TEXT_FAINT, italic=True)
    add_footer(slide, page, total)
    return slide

def slide_quote(prs, num, eyebrow, quote, attribution, page, total):
    slide = slide_blank(prs)
    add_eyebrow(slide, Inches(0.55), Inches(0.5), num, eyebrow)
    add_text(slide, Inches(0.6), Inches(2.6), Inches(12), Inches(0.5),
             '“', size=64, color=PURPLE, bold=True)
    add_text(slide, Inches(1.3), Inches(2.9), Inches(11), Inches(2.5),
             quote, size=30, color=WHITE, italic=True)
    add_text(slide, Inches(1.3), Inches(5.7), Inches(11), Inches(0.5),
             attribution, size=12, color=PURPLE_SOFT, bold=True)
    add_footer(slide, page, total)
    return slide

def slide_two_col(prs, num, eyebrow, headline, left_title, left_body, right_title, right_body, page, total):
    slide = slide_blank(prs)
    add_eyebrow(slide, Inches(0.55), Inches(0.5), num, eyebrow)
    add_text(slide, Inches(0.55), Inches(2.2), Inches(12), Inches(1.2),
             headline, size=40, bold=True, color=WHITE)
    add_purple_rule(slide, Inches(0.6), Inches(3.45), Inches(1.8), thick=Pt(3))

    col_w = Inches(5.85)
    gap   = Inches(0.45)
    for i, (title, body, tint) in enumerate([
        (left_title,  left_body,  PURPLE_INK),
        (right_title, right_body, PANEL),
    ]):
        x = Inches(0.55) + i * (col_w + gap)
        add_rect(slide, x, Inches(3.95), col_w, Inches(2.85),
                 rgb=tint if i == 0 else PANEL,
                 line_rgb=BORDER, corner=0.04)
        add_text(slide, x + Inches(0.4), Inches(4.15), col_w - Inches(0.8), Inches(0.7),
                 title, size=18, bold=True, color=WHITE)
        # purple rule under title
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                      x + Inches(0.4), Inches(4.75), Inches(0.55), Pt(2))
        line.fill.solid(); line.fill.fore_color.rgb = PURPLE
        line.line.fill.background()
        add_text(slide, x + Inches(0.4), Inches(4.95), col_w - Inches(0.8), Inches(2),
                 body, size=13, color=TEXT_DIM)
    add_footer(slide, page, total)
    return slide

def slide_closing(prs, total):
    slide = slide_blank(prs)
    add_text(slide, Inches(0.55), Inches(0.5), Inches(12), Inches(0.4),
             '— END OF PITCH', size=11, color=PURPLE_SOFT, letter_spacing=2)
    # Massive begin
    add_text(slide, Inches(0.55), Inches(1.6), Inches(12), Inches(3.5),
             'BEGIN.', size=220, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    add_purple_rule(slide, Inches(0.6), Inches(5.1), Inches(2.4), thick=Pt(4))
    add_text(slide, Inches(0.6), Inches(5.4), Inches(12), Inches(0.6),
             'Take five minutes. Open Kairo. Ask one doubt.',
             size=22, color=TEXT_DIM, italic=True)
    # Contact card
    add_rect(slide, Inches(0.55), Inches(6.3), Inches(12.2), Inches(0.65),
             rgb=PANEL, line_rgb=BORDER, corner=0.5)
    add_text(slide, Inches(0.85), Inches(6.4), Inches(6), Inches(0.5),
             'kairo-daily-edu.vercel.app', size=14, bold=True, color=PURPLE_SOFT)
    add_text(slide, Inches(7.5), Inches(6.4), Inches(5.2), Inches(0.5),
             'quro.cor@gmail.com', size=14, bold=True, color=WHITE,
             align=PP_ALIGN.RIGHT)
    add_footer(slide, total, total)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 12

    # 01 · Cover
    slide_cover(prs, TOTAL)

    # 02 · The Problem
    slide_bullets(prs, '01', 'The Problem', 'Forty million students. One textbook.', [
        ('Every student gets the same textbook.', 'Nobody gets the same brain.'),
        ('Teachers move at the speed of the syllabus.', 'Not at the speed of the student.'),
        ('Coaching costs ₹50,000+ a year.', 'Most families can\'t afford it.'),
        ('Apps exist.', 'But none remember what YOU forgot last Wednesday.'),
    ], 2, TOTAL)

    # 03 · The Big Claim
    slide_section(prs, '02', 'The Claim',
        'Kairo isn\'t a learning app. It learns you.',
        'A memory-first AI tutor that watches what you ask, where you stumble, what you replay '
        'at 1 a.m. — then tutors you back with explanations only you needed. The longer you stay, '
        'the more it becomes only yours.',
        3, TOTAL)

    # 04 · The Product
    slide_bullets(prs, '03', 'The Product', 'Eleven systems. One mind.', [
        ('Solver',         'Any doubt, any subject — full explanation in 8 seconds with sourced images.'),
        ('Kairo OS',       'The Twin. Memory engine. Tracks what you\'ve studied + forgotten + need next.'),
        ('Labs',           'Touchable 3D simulations — physics, chem, bio. Drag, pinch, learn.'),
        ('Voice tutor',    'Hold-to-speak. Hands-free, exam-night ready.'),
        ('Notebook · Concept Map · Predictor · Battle · Pomodoro · Camera · Adaptive Quiz',
                           '+ a dozen more, all wired into the same Twin.'),
    ], 4, TOTAL)

    # 05 · The Twin (deeper)
    slide_two_col(prs, '04', 'The Twin',
        'Your data lives on YOUR device.',
        'Local-first',
        'Every flashcard, formula, focus session, voice prompt — all stored in your browser\'s '
        'local storage. Nothing leaves your device unless you sync.',
        'Privacy by design',
        'When you switch devices, Kairo encrypts the whole Twin, ships it across, and wipes the '
        'cloud copy. The server is a transit lane, not a database.',
        5, TOTAL)

    # 06 · Traction stat grid
    slide_stat_grid(prs, '05', 'What\'s Shipped',
        'Built and live.', [
            ('45+',  'features in production'),
            ('27',   '3D Labs live'),
            ('8s',   'avg. Solver answer'),
            ('0',    'cost per student'),
        ], 6, TOTAL,
        note='All shipped at https://kairo-daily-edu.vercel.app — open the live status page at /status.')

    # 07 · Founder story
    slide_section(prs, '06', 'The Founder',
        'Built by a 13-year-old in Chennai.',
        'Darshan is in Class 9 at a CBSE school in India. He started Kairo because none of the '
        'apps he was given to study with knew anything about him. He spent four months teaching '
        'himself React, TypeScript, Express, Supabase, three.js, and enough Framer Motion to make '
        'diagrams move — then shipped Kairo to the web.',
        7, TOTAL)

    # 08 · Quote pull
    slide_quote(prs, '07', 'A note from the editor',
        'I wasn\'t trying to build a startup. I was trying to fix the part of school that nobody '
        'wanted to fix.',
        '— DARSHAN  ·  FOUNDER, KAIRO',
        8, TOTAL)

    # 09 · Why free
    slide_section(prs, '08', 'Why It\'s Free',
        'Kairo is free, on purpose.',
        'The students who need Kairo most are the ones whose families can\'t pay for tutors and '
        'don\'t speak the language of the coaching centres. If Kairo costs money, it stops being '
        'for them. Schools that want the whole platform can pay later — only after their students '
        'are already using it daily.',
        9, TOTAL)

    # 10 · Tech stack
    slide_bullets(prs, '09', 'Under The Hood', 'A real software stack.', [
        ('Front-end',     'React + TypeScript + Vite. Framer Motion for every transition. three.js for the 3D labs.'),
        ('Server',        'Express on Vercel. Multi-model AI race across OpenRouter + Groq + Wikipedia fallback.'),
        ('Database',      'Supabase (Postgres + auth + storage). Local-first Twin in browser localStorage.'),
        ('Email',         'Gmail SMTP via Nodemailer for OTP + transactional. Brevo/Resend fallback ready.'),
        ('Deploy',        'Auto-deploy on git push. Public status at /status with live latency + service health.'),
    ], 10, TOTAL)

    # 11 · The Ask
    slide_bullets(prs, '10', 'The Ask', 'Build with us.', [
        ('Builders',     'Designers, engineers, inventors — volunteer or equity. Remote. Async. Credit on every ship.'),
        ('Schools',      'Sign up one class. Watch what happens. We onboard, no setup fee.'),
        ('Mentors',      'A founder who\'s shipped at scale. Thirty minutes a month is enough.'),
        ('Sponsors',     'Cover Vercel Pro + Supabase Pro + OpenRouter credits (~$50/mo). Scales Kairo to 3000 daily users.'),
    ], 11, TOTAL, foot_note='email: quro.cor@gmail.com   ·   subject: I\'LL BUILD')

    # 12 · Closing
    slide_closing(prs, TOTAL)

    out = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', 'Kairo_Pitch_Deck_v2.pptx'))
    prs.save(out)
    try:
        print(f'wrote Kairo_Pitch_Deck_v2.pptx  ({TOTAL} slides)')
    except UnicodeEncodeError:
        pass


if __name__ == '__main__':
    build()
