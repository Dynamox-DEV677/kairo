"""
Generate Kairo_Walkthrough.pptx — a slide-by-slide tour of every screen
in Kairo, with every button explained in plain English.

Run:
    pip install python-pptx
    python tools/make_kairo_walkthrough.py

Output:  Kairo_Walkthrough.pptx (at repo root)
"""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ─── palette ───────────────────────────────────────────────────────────────
BG          = RGBColor(0x06, 0x06, 0x0A)
PANEL       = RGBColor(0x0C, 0x0C, 0x14)
PANEL2      = RGBColor(0x14, 0x14, 0x22)
BORDER      = RGBColor(0x22, 0x22, 0x2E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DIM    = RGBColor(0xC1, 0xC1, 0xC8)
TEXT_FAINT  = RGBColor(0x8A, 0x8A, 0x96)
PURPLE      = RGBColor(0xA7, 0x8B, 0xFA)
PURPLE_HI   = RGBColor(0x7C, 0x3A, 0xED)
PURPLE_SOFT = RGBColor(0xC4, 0xB5, 0xFD)
PURPLE_LITE = RGBColor(0xE9, 0xD5, 0xFF)
PURPLE_DEEP = RGBColor(0x5B, 0x21, 0xB6)

# 16:9 slide dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ─── helpers ───────────────────────────────────────────────────────────────
def set_solid(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    shape.line.fill.background()

def add_rect(slide, x, y, w, h, rgb=PANEL, line_rgb=None, corner=0.0):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    if corner:
        shape.adjustments[0] = corner
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    if line_rgb is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(0.75)
    shape.shadow.inherit = False
    return shape

def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font='Inter'):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def add_bullets(slide, x, y, w, h, items, *,
                size=13, bullet_color=PURPLE_SOFT,
                text_color=TEXT_DIM, bold_label=True, line_gap=8, font='Inter'):
    """Each item is (label, body). label is bold, body is dimmer."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    for i, (label, body) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(line_gap)
        # bullet dot
        r0 = p.add_run()
        r0.text = '•  '
        r0.font.name = font
        r0.font.size = Pt(size)
        r0.font.bold = True
        r0.font.color.rgb = bullet_color
        # label
        if label:
            r1 = p.add_run()
            r1.text = label
            r1.font.name = font
            r1.font.size = Pt(size)
            r1.font.bold = bold_label
            r1.font.color.rgb = WHITE
            r2 = p.add_run()
            r2.text = ' — ' + body if body else ''
            r2.font.name = font
            r2.font.size = Pt(size)
            r2.font.color.rgb = text_color
        else:
            r2 = p.add_run()
            r2.text = body
            r2.font.name = font
            r2.font.size = Pt(size)
            r2.font.color.rgb = text_color
    return tb

def add_background(slide, rgb=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb
    bg.line.fill.background()

def add_corner_glow(slide):
    """Subtle purple-tinted blocks in the corners for that Kairo vibe."""
    g1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
    g1.fill.solid()
    g1.fill.fore_color.rgb = PURPLE_DEEP
    g1.line.fill.background()
    # python-pptx can't fade real gradients via API; we approximate with a tinted oval
    # plus a translucent rect overlay via xml
    _set_alpha(g1, 0x15)  # very faint
    g2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(5), Inches(6), Inches(6))
    g2.fill.solid()
    g2.fill.fore_color.rgb = PURPLE_HI
    g2.line.fill.background()
    _set_alpha(g2, 0x10)

def _set_alpha(shape, alpha_byte):
    """alpha_byte 0x00..0xFF — but pptx uses 0..100000 in OOXML."""
    sp = shape.fill._xPr.find(qn('a:solidFill'))
    if sp is None:
        return
    clr = sp.find(qn('a:srgbClr'))
    if clr is None:
        return
    # remove any existing alpha
    for a in clr.findall(qn('a:alpha')):
        clr.remove(a)
    a = etree.SubElement(clr, qn('a:alpha'))
    val = int(alpha_byte / 255 * 100000)
    a.set('val', str(val))

def add_chip(slide, x, y, w, h, text, rgb_bg=PANEL2, rgb_border=BORDER, rgb_text=PURPLE_SOFT, size=10):
    shape = add_rect(slide, x, y, w, h, rgb=rgb_bg, line_rgb=rgb_border, corner=0.5)
    tf = shape.text_frame
    tf.margin_left = Pt(8); tf.margin_right = Pt(8)
    tf.margin_top  = Pt(2); tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text.upper()
    r.font.name = 'Inter'
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = rgb_text
    return shape

def add_header(slide, eyebrow, title, subtitle=None):
    """Standard slide header: small purple eyebrow, big title, optional subtitle."""
    add_chip(slide, Inches(0.6), Inches(0.5), Inches(2.2), Inches(0.32),
             eyebrow, rgb_text=PURPLE_SOFT)
    add_text(slide, Inches(0.6), Inches(0.95), Inches(12), Inches(0.9),
             title, size=34, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.65), Inches(12), Inches(0.5),
                 subtitle, size=14, color=TEXT_DIM)

def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
             'Kairo — Accelerate Your Academics',
             size=9, color=TEXT_FAINT)
    add_text(slide, Inches(7.3), Inches(7.05), Inches(5.5), Inches(0.3),
             f'{page_num} / {total}',
             size=9, color=TEXT_FAINT, align=PP_ALIGN.RIGHT)

# ─── slide builders ────────────────────────────────────────────────────────
def slide_blank(prs):
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    add_background(slide)
    add_corner_glow(slide)
    return slide

def slide_title(prs):
    slide = slide_blank(prs)
    # centred title
    add_text(slide, Inches(0.6), Inches(2.1), Inches(12), Inches(1),
             'Kairo',
             size=88, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(3.4), Inches(12), Inches(0.6),
             'Your AI Academic Twin',
             size=24, color=PURPLE_SOFT, align=PP_ALIGN.CENTER)
    # divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(5.6), Inches(4.4), Inches(2.1), Pt(2))
    line.fill.solid(); line.fill.fore_color.rgb = PURPLE
    line.line.fill.background()
    add_text(slide, Inches(0.6), Inches(4.7), Inches(12), Inches(0.5),
             'A complete walkthrough — every screen, every button, in plain English.',
             size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(0.4),
             'Built for Class 9–12 students in India · CBSE · ICSE · State boards',
             size=11, color=TEXT_FAINT, align=PP_ALIGN.CENTER)
    return slide

def slide_section(prs, eyebrow, title, blurb):
    slide = slide_blank(prs)
    add_chip(slide, Inches(0.6), Inches(2.5), Inches(2.2), Inches(0.36),
             eyebrow, rgb_text=PURPLE_SOFT, size=11)
    add_text(slide, Inches(0.6), Inches(3.0), Inches(12), Inches(1.2),
             title, size=54, bold=True, color=WHITE)
    add_text(slide, Inches(0.6), Inches(4.5), Inches(11), Inches(1.5),
             blurb, size=16, color=TEXT_DIM)
    return slide

def slide_screen(prs, eyebrow, title, subtitle, buttons, page, total, note=None):
    """
    buttons: list of (button_name, plain_english_explanation)
    note: optional one-line tail
    """
    slide = slide_blank(prs)
    add_header(slide, eyebrow, title, subtitle)

    # body panel
    panel = add_rect(slide, Inches(0.6), Inches(2.3),
                     Inches(12.13), Inches(4.5),
                     rgb=PANEL, line_rgb=BORDER, corner=0.04)

    add_text(slide, Inches(0.95), Inches(2.5), Inches(11.5), Inches(0.4),
             'What every button does',
             size=11, bold=True, color=PURPLE_SOFT)

    add_bullets(slide, Inches(0.95), Inches(2.95),
                Inches(11.5), Inches(3.7),
                buttons, size=12.5, line_gap=6)

    if note:
        add_text(slide, Inches(0.6), Inches(6.85), Inches(12), Inches(0.3),
                 note, size=10, color=TEXT_FAINT)

    add_footer(slide, page, total)
    return slide

# ═══════════════════════════════════════════════════════════════════════════
# CONTENT — every screen, every button explained
# ═══════════════════════════════════════════════════════════════════════════

# Each entry: (eyebrow, title, subtitle, [ (button_label, plain_english), ... ], optional_note)
SCREENS = [

# ── LANDING & AUTH ────────────────────────────────────────────────────────
('LANDING PAGE', 'The cinematic front door',
 'What new visitors see before they sign in.',
 [
  ('Get Started',          'Opens the sign-up / sign-in screen so you can create your account or log back in.'),
  ('Sign In',              'Same as Get Started — jumps straight to the email + password form for returning users.'),
  ('Try a Lab (3D demo)',  'Drops you into one of the free 3D physics simulations without signing up, just to feel the product.'),
  ('Learn More / Scroll',  'Smooth-scrolls through the marketing sections — Problem, Kairo OS, Solver, Labs, Roles, Features.'),
  ('Privacy Policy (footer)','Opens the in-app Privacy modal explaining what data we collect and where it lives.'),
  ('Terms of Service (footer)','Opens the in-app Terms modal — rules, eligibility, liability, governing law.'),
  ('Contact (footer)',     'One-click copy of the support email so you can write to us.'),
 ],
 'The landing page is server-less marketing — no account needed.'),

('SIGN IN', 'Welcome back',
 'The single email + password form for existing accounts.',
 [
  ('Email field',          'Type the email you used when you signed up.'),
  ('Password field',       'Type your password. Hidden by default.'),
  ('Eye icon',             'Tap to show or hide the password while you type — handy on mobile.'),
  ('Forgot your password?','Sends a password-reset link to the email above. You set a new password and come back to sign in.'),
  ('Sign In button',       'Logs you in. On success, takes you to the dashboard. On failure, shows the exact error from Supabase.'),
  ('Back arrow',           'Returns you to the "Choose how to sign up" screen.'),
  ('Terms & Privacy links','Each link opens the full document in a sheet — required reading on first sign-in.'),
 ],
 'Powered by Supabase auth. Sessions auto-refresh every 10 minutes.'),

('SIGN UP — PERSONAL', 'For students with no school',
 'One-screen sign-up. Skip the school code entirely.',
 [
  ('Avatar picker',        'Optional. Tap to upload a profile photo (up to 4 MB, JPG / PNG / WebP).'),
  ('Full Name',            'Your full name as it should appear inside Kairo and on shared content.'),
  ('Email',                'Your login email — also where we send OTPs and reset links.'),
  ('Create Password',      'At least 8 characters. Hidden by default; eye icon toggles visibility.'),
  ('Class / Grade',        'Optional — helps Kairo tailor explanations to your level (e.g. "9", "10 A", "Class 11").'),
  ('Board',                'Optional pill picker — CBSE, ICSE, State, IB, Other. Used by the AI to match your syllabus.'),
  ('Create my account',    'Creates the account, signs you in, and drops you into your fresh dashboard.'),
  ('Sign in instead',      'If you actually have an account, this jumps to the sign-in form.'),
 ],
 'If your email already exists, Kairo auto-signs you in with the same password.'),

('SIGN UP — JOIN A SCHOOL', '4-step wizard',
 'For students, teachers, and parents whose school is already on Kairo.',
 [
  ('Step 1 · Join Code',         'Enter the 18-character code your school admin gave you.'),
  ('Continue',                   'Validates the code with the server and shows the school logo + name on the next screen.'),
  ('Step 2 · Avatar + Account',  'Pick a photo, then enter Name, Email, Password (8+ chars).'),
  ('Step 3 · Pick role',         'Tap one of: Student, Teacher, Parent. The next screen depends on what you picked.'),
  ('Step 4 · Parent linking',    'Parents only — enter the student\'s name and the 8-char parent access code from the child\'s app.'),
  ('Create Account',             'Final submit. Logs you in with the role-appropriate dashboard.'),
  ('Back arrow on every step',   'Goes back one step without losing what you typed.'),
 ],
 None),

('SIGN UP — CREATE A SCHOOL', '2-step admin wizard',
 'For principals, owners, or coordinators starting fresh on Kairo.',
 [
  ('Step 1 · School Name', 'Type your school\'s full name as you want it shown to teachers and students.'),
  ('Step 2 · Owner Account','Your name, email, and password. You become the admin.'),
  ('Create School',        'Server creates the school, makes you admin, and signs you in immediately.'),
  ('Join Code reveal',     'Final screen shows a unique join code. Share it with teachers and students to invite them.'),
  ('Copy code',            'One-tap clipboard copy of the join code.'),
  ('Open Admin Dashboard', 'Takes you straight into the School Hub control center.'),
 ],
 'Free during early access — no payment required to start a school.'),

# ── CORE DASHBOARD ───────────────────────────────────────────────────────
("KAIRO'S SOLVER", 'Any doubt. Eight seconds.',
 'The main AI tutor — type any question, get a visual lesson with images + explanation.',
 [
  ('Ask anything box',     'Type your question (any subject, any board). Press Enter or tap Send.'),
  ('Send (▶)',             'Fires the request. Kairo writes a step-by-step explanation on the right and pulls 4–6 images on the left.'),
  ('Stop (⏹)',             'Cancels the streaming response if you typed the wrong thing.'),
  ('Suggestion chips',     'Pre-filled example questions on the empty state — tap to ask without typing.'),
  ('Image slideshow ◀ / ▶','Walk through the auto-built "picture-book" of images (Wikimedia, Pexels, Unsplash).'),
  ('Open in Kairo Labs',   'If the topic has a matching 3D lab (e.g. gravity, circuits), jumps you there with one tap.'),
  ('Mic',                  'Voice input — speak your question instead of typing.'),
  ('Model picker',         'Switches between different AI models (default is the fastest).'),
  ('New chat',             'Wipes the screen and starts fresh.'),
 ],
 None),

('KAIRO OS', 'Your AI Academic Twin',
 'The memory engine — everything Kairo knows about you, in one place.',
 [
  ('Twin tile',            'Shows your study activity, mastered topics, weak areas. Auto-updates as you learn.'),
  ('Voice',                'Hold to speak. Kairo responds in voice — fully conversational tutor.'),
  ('Insights cards',       'Daily nudges — "You\'ve forgotten 4 chem flashcards", "Time to revise Thermodynamics".'),
  ('Quick actions row',    'One-tap shortcuts: Make flashcards, Build study plan, Take quiz, Open notebook.'),
  ('Cross-device sync',    'If you sign in on another device, your whole Twin sprints over — locally-first, cloud transit only.'),
 ],
 'All Twin data lives on your device; the cloud is only used to hop devices.'),

('KAIRO LABS', '3D simulations of physics, chem, bio',
 'Touchable simulations instead of static diagrams.',
 [
  ('Lab card',             'Tap any lab tile to launch the 3D scene full-screen.'),
  ('Drag / pinch',         'Orbit and zoom the scene with your finger or mouse.'),
  ('Sliders',              'Adjust physical variables in real time (mass, voltage, temperature, etc.).'),
  ('Reset',                'Return the scene to default values without leaving the lab.'),
  ('Explain this',         'Asks Kairo Solver to explain what you\'re currently seeing in the simulation.'),
  ('Take screenshot',      'Save the current frame to your notebook as a study reference.'),
 ],
 None),

# ── DAILY STUDY ──────────────────────────────────────────────────────────
('FLASHCARDS', 'Spaced repetition done right',
 'Make, study, and review cards that come back at the right time.',
 [
  ('+ New deck',           'Create a new deck for a chapter or topic.'),
  ('AI generate',          'Paste a passage and Kairo turns it into 10–20 cards automatically.'),
  ('Study now',            'Starts a session — show card, flip, rate "Hard / Good / Easy". Next time is scheduled by SRS.'),
  ('Hard',                 'Card comes back in 10 minutes.'),
  ('Good',                 'Card comes back in 1 day (then doubles).'),
  ('Easy',                 'Card comes back in 4 days (then triples).'),
  ('Edit card',            'Tap a card to fix the question or answer.'),
  ('Delete card',          'Remove a card from the deck (asks for confirmation).'),
 ],
 None),

('NOTEBOOK', 'AI-organised study notes',
 'Type or paste anything; Kairo formats, tags, and links it.',
 [
  ('New note',             'Blank canvas. Markdown supported.'),
  ('Generate from prompt', 'Tell Kairo what chapter you\'re studying and it drafts an organised note.'),
  ('Voice note',           'Record a voice memo — Kairo transcribes and summarises.'),
  ('Save',                 'Persists the note to your Twin and the device.'),
  ('Tag chip',             'Add tags (e.g. "Physics, Class 10, Boards"). Used by search.'),
  ('Export',               'Download the note as Markdown or PDF.'),
 ],
 None),

('STUDY PLAN', 'A weekly plan that bends to your life',
 'Tell Kairo your exams and Kairo schedules everything.',
 [
  ('+ Add exam',           'Drop in an exam date and subject — Kairo back-schedules revision sessions.'),
  ('Generate plan',        'AI builds a week-by-week study plan based on your weak topics.'),
  ('Drag-to-reschedule',   'Move a session to a different day or time.'),
  ('Mark done',            'Tick off a session — feeds back into your Twin\'s memory score.'),
  ('Skip / Snooze',        'Postpone a session without losing the streak.'),
 ],
 None),

('POMODORO', '25-minute focus timer',
 'A minimal timer with break enforcement and stats.',
 [
  ('Start',                '25-minute focus block begins; phone notifications auto-silenced via the API.'),
  ('Pause',                'Stops the clock — useful if a teacher walks in.'),
  ('Skip break',           'Jumps to the next focus block without taking the 5-minute break (not recommended).'),
  ('Long break',           'Triggered automatically after every 4 blocks — 15 minutes off.'),
  ('Stats',                'Total focus time today / this week.'),
 ],
 None),

('FOCUS MODE', 'Hide everything except what you\'re studying',
 'Fullscreen single-task UI for deep work.',
 [
  ('Pick a target',        'Choose a subject + topic for this focus session.'),
  ('Start session',        'Hides side panels, mutes notifications, dims the rest of the UI.'),
  ('Exit',                 'Restores the full dashboard. Asks "Was this useful?" to feed your Twin.'),
 ],
 None),

('PANIC MODE', 'Exam tomorrow. Don\'t panic.',
 'A calm, last-minute revision guide built for night-before-the-exam moments.',
 [
  ('Subject picker',       'Pick the subject you\'re panicking about.'),
  ('Top 10 must-knows',    'Kairo surfaces the 10 highest-yield concepts based on your weak areas.'),
  ('5-minute summaries',   'Each concept is a 5-min revision card with example + answer.'),
  ('Quick quiz',           'A 10-question rapid quiz to lock memory in before sleep.'),
 ],
 None),

('MEMORY BRAIN', 'The Twin\'s memory graph',
 'See everything Kairo remembers about how you learn.',
 [
  ('Strong / Weak chips',  'Topics colour-coded by mastery — light = strong, dark = weak.'),
  ('Replay a memory',      'Tap a topic to see when you last studied it, your last score, and what to do next.'),
  ('Forget this',          'Manually mark a topic as forgotten — schedules instant revision.'),
  ('Export brain',         'Download a JSON snapshot of your full Twin.'),
 ],
 None),

('CONCEPT MAP', 'How everything connects',
 'A draggable graph linking concepts, formulas, and examples.',
 [
  ('Node',                 'Tap a concept to expand its definition + related nodes.'),
  ('Drag',                 'Reposition the graph for a clearer view.'),
  ('Zoom',                 'Pinch or wheel-scroll to zoom in / out.'),
  ('Search bar',           'Jump directly to any concept on the graph.'),
  ('Add link',             'Manually connect two concepts you discovered are related.'),
 ],
 None),

('FORMULA SHEET', 'Every formula you\'ve seen',
 'Searchable, swipeable cheat sheet auto-built from your studied topics.',
 [
  ('Search',               'Type "kinematics" to filter to relevant formulas.'),
  ('Tap a formula',        'Opens a full derivation + example problem.'),
  ('Pin',                  'Pin to the top of the sheet for fast access during revision.'),
  ('+ Add formula',        'Save a formula you ran into while studying.'),
 ],
 None),

# ── EXAM PREP ────────────────────────────────────────────────────────────
('ADAPTIVE QUIZ', 'Quizzes that change as you go',
 'Difficulty scales up or down based on your last 3 answers.',
 [
  ('Start quiz',           'Begins a 20-question quiz seeded from your weak topics.'),
  ('Submit answer',        'Locks your answer and tells you correct / wrong with the model explanation.'),
  ('Hint',                 'Asks Kairo for a small hint without revealing the answer (counts against the score).'),
  ('Skip',                 'Move on without answering — comes back at the end.'),
  ('End early',            'Stop the quiz and see partial results.'),
 ],
 None),

('BATTLE MODE', '1-v-1 quiz duels',
 'Pick a friend\'s code or random match — same questions, fastest correct answer wins.',
 [
  ('Join battle',          'Enter the 6-char code your friend shared.'),
  ('Random match',         'Auto-match against another Kairo user at your level.'),
  ('Start',                'Begins the 10-question duel — 15 seconds per question.'),
  ('Lock answer',          'Submits your pick. First correct answer takes the point.'),
  ('Rematch',              'After the battle, instantly start another with the same opponent.'),
 ],
 None),

('EXAM PREDICTOR', 'What will likely appear',
 'AI predicts the highest-probability questions for your upcoming exam.',
 [
  ('Pick exam',            'Choose the board + class + subject + exam type.'),
  ('Generate prediction',  'Returns 15–25 likely questions ranked by probability.'),
  ('See past appearances', 'Tap a question to see when this exact concept appeared in past years.'),
  ('Practice this',        'Sends the predicted question to Kairo Solver for a full walkthrough.'),
 ],
 None),

('PERFORMANCE PREDICTOR', 'Your likely score',
 'Honest estimate of where you\'ll land based on your current Twin.',
 [
  ('Estimate score',       'Returns an expected percentage with a confidence range.'),
  ('What\'s holding you back','Top 3 weak topics that are pulling the prediction down.'),
  ('Improve this',         'Generates a 7-day plan focused on each weak topic.'),
 ],
 None),

('MISTAKE ANALYSIS', 'What you keep getting wrong',
 'Patterns Kairo notices across your wrong answers.',
 [
  ('Wrong answers list',   'Every wrong attempt across quizzes and battles, grouped by topic.'),
  ('Pattern card',         'Plain-English explanation of WHY you tend to miss this kind of question.'),
  ('Drill these',          'Generates 10 practice questions targeting the same mistake type.'),
 ],
 None),

('QUESTION PAPER', 'Mock paper generator',
 'A full mock exam paper in your board\'s format.',
 [
  ('Pick board + class',   'CBSE / ICSE / State + 9–12.'),
  ('Choose subject',       'Math, Physics, Chem, Bio, English, etc.'),
  ('Generate',             'Builds a complete paper — section A/B/C, marks distribution, time limit.'),
  ('Start timer',          'Begins a timed attempt — clock counts down in the corner.'),
  ('Submit',               'Auto-grades MCQs; sends written answers to the Essay Grader.'),
  ('Download PDF',         'Save the paper to print or share with classmates.'),
 ],
 None),

('REVISION SIMULATOR', 'Simulate the exam day',
 'Practice the full exam experience — pressure, timing, sequencing.',
 [
  ('Begin simulation',     'Locks the UI into exam mode — no nav, no notifications.'),
  ('Question nav',         'Jump between questions like in a real online paper.'),
  ('Flag for review',      'Mark a question to revisit before submitting.'),
  ('Submit paper',         'Ends the simulation and shows score + per-question breakdown.'),
 ],
 None),

# ── OTHER TOOLS ──────────────────────────────────────────────────────────
('VOICE TUTOR', 'Talk to Kairo',
 'Conversational voice tutor for hands-free learning.',
 [
  ('Hold to speak',        'Records your question while pressed.'),
  ('Release',              'Sends the audio. Kairo replies in voice + on-screen text.'),
  ('Interrupt',            'Tap once to cut off Kairo mid-sentence (useful for clarifications).'),
  ('Save to notebook',     'Stores the conversation as a note.'),
 ],
 None),

('CAMERA STUDY', 'Point at a textbook',
 'Snap a page, get an instant explanation.',
 [
  ('Open camera',          'Asks browser permission, then shows the live viewfinder.'),
  ('Capture',              'Takes a photo of the page or problem.'),
  ('Explain',              'Sends the photo to Kairo — returns a step-by-step solution.'),
  ('Save',                 'Adds the photo + explanation to your notebook.'),
 ],
 None),

('ESSAY GRADER', 'AI English-teacher feedback',
 'Paste an essay, get a grade plus inline comments.',
 [
  ('Paste essay',          'Drop your essay into the text box.'),
  ('Grade',                'Returns a score (out of 10), strengths, weaknesses, and rewrite suggestions.'),
  ('Inline comments',      'Hover over highlighted sentences to see what to improve.'),
  ('Apply rewrite',        'Accept Kairo\'s suggested rewrite for a paragraph.'),
 ],
 None),

('WRITING TOOLS', 'Brainstorm, outline, polish',
 'AI helpers for any writing assignment.',
 [
  ('Brainstorm ideas',     'Generates 8–10 angles for an essay prompt.'),
  ('Outline',              'Turns your chosen angle into a 3-part outline.'),
  ('Expand paragraph',     'Takes a sentence and fleshes it into a full paragraph.'),
  ('Tone-shift',           'Rewrites in formal / casual / academic tone.'),
 ],
 None),

('KNOWLEDGE GRAPH', 'Your whole syllabus as a map',
 'The full subject graph — like a metro map for your textbook.',
 [
  ('Subject pill',         'Filter to one subject at a time.'),
  ('Tap a node',           'See definition + linked sub-topics.'),
  ('Path-find',            'Pick two nodes; Kairo shows the shortest learning path between them.'),
 ],
 None),

('GAMIFICATION', 'XP, streaks, badges',
 'Optional layer that rewards consistent studying.',
 [
  ('Daily streak',         'Counts consecutive study days — visible at the top of the dashboard.'),
  ('Claim XP',              'Tap when you finish a session to log XP.'),
  ('Badges tab',           'Shows all the achievements unlocked so far.'),
  ('Leaderboard (optional)','Compare with classmates in the same school — opt-in only.'),
 ],
 None),

# ── TEACHER TOOLS ────────────────────────────────────────────────────────
('LESSON PLAN', 'AI lesson planner',
 'For teachers — instantly draft a lesson plan for any chapter.',
 [
  ('Pick chapter',         'Subject + grade + chapter.'),
  ('Generate plan',        '60-min lesson with objectives, activities, assessment, homework.'),
  ('Edit inline',          'Tweak any section by typing.'),
  ('Export to PDF',        'Print-ready lesson sheet.'),
 ],
 None),

('TEACHER ASSISTANT', 'AI co-pilot for teachers',
 'A quick assistant for grading, doubts, and class prep.',
 [
  ('Ask anything',         'A teacher-tuned version of the Solver — gets pedagogy-aware answers.'),
  ('Bulk-grade',           'Upload a folder of student answers; Kairo grades each one.'),
  ('Talking points',       'Generate discussion points for tomorrow\'s class on a given topic.'),
 ],
 None),

('EXPLAIN MISTAKE', 'For teachers reviewing student work',
 'See WHY a student keeps making a specific mistake.',
 [
  ('Pick student',         'Choose a student from your class roster.'),
  ('See common mistakes',  'Top 5 wrong-answer patterns from their Kairo activity.'),
  ('Draft a message',      'Generate a kind, encouraging message to the student about one of those patterns.'),
 ],
 None),

# ── ADMIN / SCHOOL TOOLS ─────────────────────────────────────────────────
('SCHOOL HUB', 'The admin control center',
 'Run an entire school from one screen.',
 [
  ('Active students',      'Live count of students using Kairo this week.'),
  ('Invite teacher',       'Send the join code to a teacher; opens their inbox.'),
  ('Manage classes',       'Add / remove classes, assign teachers, view rosters.'),
  ('Reports',              'Weekly digest — attendance, top performers, struggling students.'),
  ('Settings',             'School logo, name, billing (when paid plans launch).'),
 ],
 None),

('OPS DASHBOARD', 'System health at a glance',
 'For admins — make sure everything is running.',
 [
  ('Live student count',   'How many students are inside Kairo right now.'),
  ('Error rate',           'Server errors in the last 24h.'),
  ('Slow queries',         'Pages that took >2s to load — investigate.'),
 ],
 None),

('ANALYTICS', 'Numbers + charts',
 'Class- and student-level analytics for teachers and admins.',
 [
  ('Time-range picker',    'Today, This week, Month, All time.'),
  ('Class chart',          'Avg quiz score per class.'),
  ('Student drill-down',   'Click a class to see per-student scores.'),
  ('Export CSV',           'Download the raw numbers.'),
 ],
 None),

('ATTENDANCE', 'Daily roll-call',
 'Quick, swipe-based attendance for teachers.',
 [
  ('Swipe right',          'Mark present.'),
  ('Swipe left',           'Mark absent.'),
  ('Long-press',           'Mark late or excused.'),
  ('Submit',               'Saves today\'s attendance to the school record.'),
 ],
 None),

('ANNOUNCEMENT', 'Talk to the whole school',
 'Push an announcement to students, parents, or teachers.',
 [
  ('Compose',              'Type your message; Markdown supported.'),
  ('Audience picker',      'All / one class / teachers only / parents only.'),
  ('Send',                 'Pushes the announcement — students see it on their dashboard, parents in the parent app.'),
 ],
 None),

('FEE REMINDER', 'Send polite fee nudges',
 'For admins — remind parents about pending fees.',
 [
  ('Pick parent',          'Search by student name.'),
  ('Draft reminder',       'AI writes a polite, culturally-aware message.'),
  ('Send',                 'Delivers via email and in-app notification.'),
 ],
 None),

('ADMISSION BOT', 'Parent enquiries on autopilot',
 'A chatbot for the school website that answers admissions FAQs.',
 [
  ('Train',                'Upload your prospectus / handbook; bot learns your school\'s answers.'),
  ('Preview',              'Test the bot before going live.'),
  ('Embed',                'Copy a one-line embed code for your school website.'),
 ],
 None),

('TIMETABLE', 'Weekly class schedule',
 'Drag-and-drop timetable builder.',
 [
  ('Drag class',           'Move a subject to a different time slot.'),
  ('Add period',           'Insert a new period for a teacher.'),
  ('Auto-fill',            'AI fills empty slots based on class needs.'),
  ('Publish',              'Pushes the timetable live to every student and teacher.'),
 ],
 None),

('PARENT DASHBOARD', 'What your child is doing on Kairo',
 'For parents linked to a student account.',
 [
  ('Today\'s study time',  'How long your child studied today.'),
  ('Subjects covered',     'Topics they worked on.'),
  ('Quiz scores',          'Latest scores from any Kairo quiz.'),
  ('Message child',        'Send an encouraging note that pops up on their dashboard.'),
 ],
 None),

('PARENT MESSAGE', 'Talk to the teacher',
 'Direct, private message thread with a teacher.',
 [
  ('New thread',           'Start a conversation about your child.'),
  ('Send',                 'Pushes the message; teacher gets a notification.'),
  ('Translate',            'Auto-translate the teacher\'s reply if it\'s in a different language.'),
 ],
 None),

# ── SETTINGS & FLOWS ────────────────────────────────────────────────────
('SETTINGS', 'Your account & app',
 'Everything you can change about Kairo.',
 [
  ('Profile',              'Change your name, avatar, class, board.'),
  ('Theme toggle',         'Dark / light mode.'),
  ('Reset Passcode',       'Opens the 5-step OTP flow to reset your 6-digit Kairo OS device PIN.'),
  ('Export Twin',          'Download a JSON of your Twin data — back up everything to your computer.'),
  ('Import Twin',          'Load a Twin JSON onto this device — useful for moving devices manually.'),
  ('Terms & Conditions',   'Opens the Terms modal.'),
  ('Privacy Policy',       'Opens the Privacy modal.'),
  ('Delete account',       'Permanently delete your account and server data within 30 days.'),
  ('Logout',               'Signs you out on this device only — other devices stay signed in.'),
 ],
 None),

('RESET PASSCODE', '5-step mobile-first flow',
 'For when you forgot your 6-digit Kairo OS device passcode.',
 [
  ('Step 1 · Email',       'Type the email on your account; tap Continue. Server sends a 6-digit code.'),
  ('Step 2 · Enter OTP',   '6-cell input. Paste-friendly. Live "Resend in 30s" countdown. Errors shake the field.'),
  ('Resend code',          'Available after 30 seconds. Up to 4 sends per 10 minutes.'),
  ('Step 3 · New PIN',     'Pick a new 6-digit PIN. Live strength meter (weak / good / strong).'),
  ('Step 4 · Confirm PIN', 'Re-type the same 6 digits. Mismatch shakes the keypad.'),
  ('Step 5 · Success',     'Confetti burst + "All set". Tap Done to return to settings.'),
  ('Back arrow on every step','Goes back one step without losing your progress.'),
 ],
 'OTP is hashed (SHA-256 + per-email salt), stored 10 min max, destroyed on first verify.'),

# ── MOBILE ──────────────────────────────────────────────────────────────
('MOBILE EXPERIENCE', 'Native-feel on small screens',
 'Kairo OS, dock, and gestures rebuilt for phones.',
 [
  ('Splash screen',        '3-second boot animation on first open per session.'),
  ('Floating dock',        'Glass-style bottom dock with the 5 most-used features.'),
  ('Pull to refresh',      'Refresh the Twin from the dashboard.'),
  ('Swipe nav',            'Swipe between screens within a feature (e.g. flashcards).'),
  ('Safe-area handling',   'All screens respect iOS / Android safe-area insets — never under the notch.'),
 ],
 None),

]

# ── TECH STACK + CLOSING ─────────────────────────────────────────────────
def slide_tech(prs, page, total):
    slide = slide_blank(prs)
    add_header(slide, 'UNDER THE HOOD', 'What powers Kairo',
               'A modern stack chosen for speed, privacy, and developer happiness.')
    # 2-column grid
    items_left = [
      ('React + TypeScript + Vite','Front-end app — fast HMR, type-safe everywhere.'),
      ('Framer Motion',            'All page transitions, micro-animations.'),
      ('Three.js / R3F',           'The 3D physics labs.'),
      ('Tailwind-free CSS-in-JS',  'Strict monochrome palette — hand-written for the brand.'),
      ('localStorage Twin',        'Your data lives on YOUR device first.'),
    ]
    items_right = [
      ('Express.js + Node',        'API server — auth, OTP, school management.'),
      ('Supabase',                 'Postgres + auth + storage.'),
      ('OpenRouter + Groq',        'AI inference for Solver, Voice Tutor, every AI feature.'),
      ('Nodemailer (Gmail SMTP)',  'Transactional email — OTPs, reset links.'),
      ('Vercel',                   'Hosting + edge functions.'),
    ]
    add_bullets(slide, Inches(0.6), Inches(2.5),
                Inches(6.0), Inches(4.2),
                items_left, size=12.5, line_gap=8)
    add_bullets(slide, Inches(7.0), Inches(2.5),
                Inches(6.0), Inches(4.2),
                items_right, size=12.5, line_gap=8)
    add_footer(slide, page, total)
    return slide

def slide_closing(prs, page, total):
    slide = slide_blank(prs)
    add_text(slide, Inches(0.6), Inches(2.5), Inches(12), Inches(1.2),
             'Built for you.',
             size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(3.7), Inches(12), Inches(0.5),
             'Every screen in Kairo was designed to help one student do one thing better:',
             size=15, color=TEXT_DIM, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(4.2), Inches(12), Inches(0.5),
             'understand.',
             size=22, bold=True, color=PURPLE_SOFT, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(0.6), Inches(5.7), Inches(12), Inches(0.4),
             'quro.cor@gmail.com  ·  kairo.app  ·  Accelerate Your Academics',
             size=11, color=TEXT_FAINT, align=PP_ALIGN.CENTER)
    add_footer(slide, page, total)
    return slide

# ═══════════════════════════════════════════════════════════════════════════
def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Compute total page count for footers
    total = 1 + 1 + 1 + len(SCREENS) + 1 + 1
    # title + overview + section + screens + tech + closing

    # Slide 1 — Title
    slide_title(prs)

    # Slide 2 — What Kairo is
    s = slide_blank(prs)
    add_header(s, 'OVERVIEW', 'What is Kairo?',
               'Kairo is an AI tutor + memory engine + study OS — built for Indian classrooms.')
    items = [
      ('AI tutor',           'Ask any doubt, get a step-by-step explanation with images and video — in 8 seconds.'),
      ('Memory engine',      'Kairo remembers what you\'ve studied, what you\'ve forgotten, and what to do next.'),
      ('Study OS',           'Flashcards, study plans, timetables, focus mode, panic mode — everything you need, one app.'),
      ('Privacy first',      'Your data lives on YOUR device. Cloud is only used to hop between devices.'),
      ('For everyone',       'Students, teachers, parents, and entire schools — each with their own dashboard.'),
    ]
    add_bullets(s, Inches(0.6), Inches(2.5),
                Inches(12.13), Inches(4.5),
                items, size=14, line_gap=10)
    add_footer(s, 2, total)

    # Slide 3 — section divider
    slide_section(prs,
        'WHAT FOLLOWS',
        'Every screen. Every button.',
        'The next slides walk through each surface in Kairo and explain — in plain '
        'English — exactly what every button does. No jargon, no marketing.')
    # set the footer for the section divider
    slide = prs.slides[2]
    add_footer(slide, 3, total)

    # Slides 4..N — every screen
    page = 3
    for i, (eyebrow, title, subtitle, buttons, note) in enumerate(SCREENS, start=4):
        page = i
        slide_screen(prs, eyebrow, title, subtitle, buttons, page, total, note)

    # Tech stack + closing
    slide_tech(prs, page + 1, total)
    slide_closing(prs, page + 2, total)

    out = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', 'Kairo_Walkthrough.pptx'))
    prs.save(out)
    try:
        print(f'wrote Kairo_Walkthrough.pptx ({page + 2} slides)')
    except UnicodeEncodeError:
        pass


if __name__ == '__main__':
    build()
